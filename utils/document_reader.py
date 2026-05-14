import PyPDF2
from pptx import Presentation
from PIL import Image
import pytesseract
import io

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def clean_text(text):
    text = text.replace("\n", " ")
    text = " ".join(text.split())

    sentences = text.split(".")
    cleaned = []

    for s in sentences:
        if len(s.strip()) > 40:
            cleaned.append(s.strip())

    return ". ".join(cleaned)


def extract_text_from_pdf(uploaded_file):
    reader = PyPDF2.PdfReader(uploaded_file)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return clean_text(text)


def extract_text_from_ppt(uploaded_file):
    presentation = Presentation(uploaded_file)
    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "

    return clean_text(text)


def extract_text_from_image(uploaded_file):
    image = Image.open(io.BytesIO(uploaded_file.read()))
    text = pytesseract.image_to_string(image)
    return clean_text(text)